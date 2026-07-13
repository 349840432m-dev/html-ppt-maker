#!/usr/bin/env python3
"""Generate an editable, simplified .pptx fallback from deck-plan.json.

This intentionally does not preserve HTML visuals. Use export_html_matched_pptx.mjs
for the recommended viewing copy. This script is only for users who explicitly
prioritize editable text over HTML/PPT visual parity.

Usage:
    python3 generate_pptx_from_plan.py deck-plan.json deck.pptx
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt
except ImportError:
    print(
        "python-pptx is required. Install with: pip install python-pptx\n"
        "(or deliver HTML + deck-plan + this command per 04-PPT下载与转场.md)",
        file=sys.stderr,
    )
    raise SystemExit(2)

SLIDE_W = Emu(12192000)  # 13.333in * 914400
SLIDE_H = Emu(6858000)   # 7.5in

INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT = RGBColor(0x55, 0x58, 0x5C)
RED = RGBColor(0xD0, 0x2A, 0x1E)


def add_textbox(slide, left_in, top_in, width_in, height_in):
    inch = 914400
    return slide.shapes.add_textbox(
        Emu(int(left_in * inch)), Emu(int(top_in * inch)),
        Emu(int(width_in * inch)), Emu(int(height_in * inch)),
    )


def set_text(box, text, size, bold=False, color=INK):
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def build_slide(prs, spec, index, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    is_hero = str(spec.get("type", "")) in {"cover", "closing"}

    title = str(spec.get("title", "")).strip() or f"Slide {index + 1}"
    takeaway = str(spec.get("takeaway", "")).strip()
    body = [str(item) for item in spec.get("body", []) if str(item).strip()]
    evidence = [str(item) for item in spec.get("evidence", []) if str(item).strip()]

    if is_hero:
        set_text(add_textbox(slide, 0.9, 2.2, 11.5, 1.6), title, 44, bold=True)
        if takeaway:
            set_text(add_textbox(slide, 0.9, 3.9, 11.0, 1.0), takeaway, 20, color=INK_SOFT)
    else:
        set_text(add_textbox(slide, 0.7, 0.5, 12.0, 1.0), title, 30, bold=True)
        if takeaway:
            set_text(add_textbox(slide, 0.7, 1.45, 12.0, 0.7), takeaway, 16, bold=True, color=RED)
        if body:
            box = add_textbox(slide, 0.7, 2.4, 12.0, 3.8)
            frame = box.text_frame
            frame.word_wrap = True
            for i, item in enumerate(body):
                para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                run = para.add_run()
                run.text = f"· {item}"
                run.font.size = Pt(17)
                run.font.color.rgb = INK
                para.space_after = Pt(10)

    footer_parts = []
    if evidence:
        footer_parts.append("来源: " + "; ".join(evidence))
    footer_parts.append(f"{index + 1} / {total}")
    set_text(add_textbox(slide, 0.7, 6.9, 12.0, 0.4), "    ".join(footer_parts), 10, color=INK_SOFT)

    notes = str(spec.get("speaker_notes", "")).strip()
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def main() -> int:
    if len(sys.argv) != 3:
        print("Usage: generate_pptx_from_plan.py <deck-plan.json> <output.pptx>", file=sys.stderr)
        return 2

    plan_path, out_path = Path(sys.argv[1]), Path(sys.argv[2])
    data = json.loads(plan_path.read_text(encoding="utf-8"))
    slides = data.get("slides", [])
    if not slides:
        print(f"{plan_path}: no slides found", file=sys.stderr)
        return 1
    if any("visible_steps" in slide for slide in slides if isinstance(slide, dict)):
        print(
            "[FAIL] storyboard-plan visual states cannot be represented by the editable fallback. "
            "Use export_html_matched_pptx.mjs with --plan and --storyboard-plan.",
            file=sys.stderr,
        )
        return 1

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for i, spec in enumerate(slides):
        build_slide(prs, spec, i, len(slides))

    prs.save(out_path)
    print(f"[WARN] wrote editable simplified fallback: {out_path} with {len(slides)} slide(s)")
    print("[WARN] this file does not preserve HTML layout; do not present it as HTML-matched")
    print("Next: python3 scripts/apply_ppt_transitions.py "
          f"{out_path} --plan {plan_path} --output {out_path.with_suffix('')}-with-transitions.pptx")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
